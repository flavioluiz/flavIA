"""Tests for the Office document converter."""

from pathlib import Path

import pytest

from flavia.content.converters import OfficeConverter, converter_registry


class TestOfficeConverterBasics:
    """Basic tests for OfficeConverter class."""

    def test_supported_extensions(self):
        """OfficeConverter supports expected extensions."""
        converter = OfficeConverter()
        expected = {
            ".docx", ".xlsx", ".pptx",  # Modern Office
            ".doc", ".xls", ".ppt",      # Legacy Office
            ".odt", ".ods", ".odp",      # OpenDocument
        }
        assert converter.supported_extensions == expected

    def test_can_handle_docx(self, tmp_path):
        """OfficeConverter can handle .docx files."""
        converter = OfficeConverter()
        docx_file = tmp_path / "test.docx"
        docx_file.write_bytes(b"PK")  # ZIP signature
        assert converter.can_handle(docx_file) is True

    def test_can_handle_xlsx(self, tmp_path):
        """OfficeConverter can handle .xlsx files."""
        converter = OfficeConverter()
        xlsx_file = tmp_path / "test.xlsx"
        xlsx_file.write_bytes(b"PK")
        assert converter.can_handle(xlsx_file) is True

    def test_can_handle_pptx(self, tmp_path):
        """OfficeConverter can handle .pptx files."""
        converter = OfficeConverter()
        pptx_file = tmp_path / "test.pptx"
        pptx_file.write_bytes(b"PK")
        assert converter.can_handle(pptx_file) is True

    def test_cannot_handle_unsupported(self, tmp_path):
        """OfficeConverter does not handle unsupported extensions."""
        converter = OfficeConverter()
        txt_file = tmp_path / "test.txt"
        txt_file.write_text("hello")
        assert converter.can_handle(txt_file) is False

    def test_requires_dependencies(self):
        """OfficeConverter declares required dependencies."""
        converter = OfficeConverter()
        assert "python-docx" in converter.requires_dependencies
        assert "openpyxl" in converter.requires_dependencies
        assert "python-pptx" in converter.requires_dependencies

    def test_dependency_import_map(self):
        """OfficeConverter maps package names to import names."""
        converter = OfficeConverter()
        assert converter.dependency_import_map.get("python-docx") == "docx"
        assert converter.dependency_import_map.get("openpyxl") == "openpyxl"
        assert converter.dependency_import_map.get("python-pptx") == "pptx"


class TestOfficeConverterRegistration:
    """Tests for OfficeConverter registry integration."""

    def test_office_converter_registered(self):
        """OfficeConverter is auto-registered in the global registry."""
        converter = converter_registry.get_for_extension(".docx")
        assert converter is not None
        assert isinstance(converter, OfficeConverter)

    def test_xlsx_registered(self):
        """OfficeConverter is registered for .xlsx."""
        converter = converter_registry.get_for_extension(".xlsx")
        assert converter is not None
        assert isinstance(converter, OfficeConverter)

    def test_pptx_registered(self):
        """OfficeConverter is registered for .pptx."""
        converter = converter_registry.get_for_extension(".pptx")
        assert converter is not None
        assert isinstance(converter, OfficeConverter)

    def test_legacy_doc_registered(self):
        """OfficeConverter is registered for legacy .doc."""
        converter = converter_registry.get_for_extension(".doc")
        assert converter is not None
        assert isinstance(converter, OfficeConverter)


class TestOfficeConverterCheckDependencies:
    """Tests for dependency checking."""

    def test_check_dependencies_returns_tuple(self):
        """check_dependencies returns a tuple."""
        converter = OfficeConverter()
        result = converter.check_dependencies()
        assert isinstance(result, tuple)
        assert len(result) == 2

    def test_check_dependencies_all_installed(self):
        """When all deps installed, returns (True, [])."""
        docx = pytest.importorskip("docx")
        openpyxl = pytest.importorskip("openpyxl")
        pptx = pytest.importorskip("pptx")

        converter = OfficeConverter()
        installed, missing = converter.check_dependencies()
        assert installed is True
        assert missing == []


class TestDocxExtraction:
    """Tests for Word document extraction."""

    @pytest.fixture
    def sample_docx(self, tmp_path):
        """Create a sample Word document for testing."""
        docx = pytest.importorskip("docx")
        from docx import Document
        from docx.shared import Pt

        doc = Document()
        doc.add_heading("Test Document", 0)
        doc.add_heading("Introduction", level=1)
        doc.add_paragraph("This is the first paragraph.")
        doc.add_heading("Details", level=2)
        doc.add_paragraph("This is another paragraph with details.")

        # Add a table
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Header 1"
        table.cell(0, 1).text = "Header 2"
        table.cell(1, 0).text = "Value 1"
        table.cell(1, 1).text = "Value 2"

        doc_path = tmp_path / "test.docx"
        doc.save(doc_path)
        return doc_path

    def test_extract_from_docx(self, sample_docx):
        """Extract text from a Word document."""
        converter = OfficeConverter()
        text = converter.extract_text(sample_docx)

        assert text is not None
        assert "Test Document" in text or "test" in text.lower()
        assert "Introduction" in text
        assert "first paragraph" in text

    def test_convert_docx_to_markdown(self, sample_docx, tmp_path):
        """Convert Word document to markdown file."""
        converter = OfficeConverter()
        output_dir = tmp_path / ".converted"

        result = converter.convert(sample_docx, output_dir, output_format="md")

        assert result is not None
        assert result.exists()
        assert result.suffix == ".md"

        content = result.read_text(encoding="utf-8")
        assert "Introduction" in content

    def test_docx_heading_detection(self, sample_docx):
        """Headings in Word document are converted to markdown headings."""
        converter = OfficeConverter()
        text = converter.extract_text(sample_docx)

        assert text is not None
        # Should have markdown heading markers
        assert "## Introduction" in text or "# " in text

    def test_docx_table_to_markdown(self, sample_docx):
        """Tables in Word document are converted to markdown tables."""
        converter = OfficeConverter()
        text = converter.extract_text(sample_docx)

        assert text is not None
        # Check for markdown table format
        assert "|" in text
        assert "Header 1" in text
        assert "Value 1" in text


class TestXlsxExtraction:
    """Tests for Excel spreadsheet extraction."""

    @pytest.fixture
    def sample_xlsx(self, tmp_path):
        """Create a sample Excel spreadsheet for testing."""
        openpyxl = pytest.importorskip("openpyxl")
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Data"

        # Add header row
        ws["A1"] = "Name"
        ws["B1"] = "Value"
        ws["C1"] = "Category"

        # Add data rows
        ws["A2"] = "Item 1"
        ws["B2"] = 100
        ws["C2"] = "A"
        ws["A3"] = "Item 2"
        ws["B3"] = 200
        ws["C3"] = "B"

        # Add second sheet
        ws2 = wb.create_sheet("Summary")
        ws2["A1"] = "Total"
        ws2["B1"] = 300

        xlsx_path = tmp_path / "test.xlsx"
        wb.save(xlsx_path)
        return xlsx_path

    def test_extract_from_xlsx(self, sample_xlsx):
        """Extract text from an Excel spreadsheet."""
        converter = OfficeConverter()
        text = converter.extract_text(sample_xlsx)

        assert text is not None
        assert "Name" in text
        assert "Item 1" in text
        assert "100" in text

    def test_xlsx_sheet_headers(self, sample_xlsx):
        """Each sheet is labeled with its name."""
        converter = OfficeConverter()
        text = converter.extract_text(sample_xlsx)

        assert text is not None
        assert "## Sheet: Data" in text
        assert "## Sheet: Summary" in text

    def test_xlsx_markdown_table_format(self, sample_xlsx):
        """Excel data is formatted as markdown tables."""
        converter = OfficeConverter()
        text = converter.extract_text(sample_xlsx)

        assert text is not None
        # Check for markdown table format
        assert "|" in text
        assert "---" in text  # Header separator

    def test_convert_xlsx_to_markdown(self, sample_xlsx, tmp_path):
        """Convert Excel spreadsheet to markdown file."""
        converter = OfficeConverter()
        output_dir = tmp_path / ".converted"

        result = converter.convert(sample_xlsx, output_dir, output_format="md")

        assert result is not None
        assert result.exists()
        assert result.suffix == ".md"

    def test_xlsx_escapes_pipe_character_in_header(self, tmp_path):
        """Header cells with pipes are escaped to keep markdown table structure valid."""
        pytest.importorskip("openpyxl")
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws["A1"] = "Name|Alias"
        ws["B1"] = "Value"
        ws["A2"] = "Item 1"
        ws["B2"] = 10

        xlsx_path = tmp_path / "pipes.xlsx"
        wb.save(xlsx_path)

        text = OfficeConverter().extract_text(xlsx_path)
        assert text is not None
        assert "| Name\\|Alias | Value |" in text


class TestPptxExtraction:
    """Tests for PowerPoint presentation extraction."""

    @pytest.fixture
    def sample_pptx(self, tmp_path):
        """Create a sample PowerPoint presentation for testing."""
        pptx = pytest.importorskip("pptx")
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()

        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Test Presentation"
        subtitle.text = "A sample presentation"

        # Add content slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide2.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "Key Points"
        tf = body_shape.text_frame
        tf.text = "First bullet point"
        p = tf.add_paragraph()
        p.text = "Second bullet point"
        p.level = 1

        pptx_path = tmp_path / "test.pptx"
        prs.save(pptx_path)
        return pptx_path

    def test_extract_from_pptx(self, sample_pptx):
        """Extract text from a PowerPoint presentation."""
        converter = OfficeConverter()
        text = converter.extract_text(sample_pptx)

        assert text is not None
        assert "Test Presentation" in text
        assert "Key Points" in text

    def test_pptx_slide_numbers(self, sample_pptx):
        """Slides are numbered in the output."""
        converter = OfficeConverter()
        text = converter.extract_text(sample_pptx)

        assert text is not None
        assert "## Slide 1" in text
        assert "## Slide 2" in text

    def test_pptx_bullet_points(self, sample_pptx):
        """Bullet points are converted to markdown lists."""
        converter = OfficeConverter()
        text = converter.extract_text(sample_pptx)

        assert text is not None
        assert "First bullet point" in text
        # Should have list markers
        assert "- " in text or "* " in text

    def test_convert_pptx_to_markdown(self, sample_pptx, tmp_path):
        """Convert PowerPoint presentation to markdown file."""
        converter = OfficeConverter()
        output_dir = tmp_path / ".converted"

        result = converter.convert(sample_pptx, output_dir, output_format="md")

        assert result is not None
        assert result.exists()
        assert result.suffix == ".md"


class TestLegacyFormatFallback:
    """Tests for legacy format handling."""

    def test_find_libreoffice(self):
        """_find_libreoffice returns a path or None."""
        converter = OfficeConverter()
        result = converter._find_libreoffice()
        # Result should be a string path or None
        assert result is None or isinstance(result, str)

    def test_legacy_extension_mapping(self):
        """Legacy extensions map to modern equivalents."""
        converter = OfficeConverter()
        assert converter._legacy_to_modern[".doc"] == ".docx"
        assert converter._legacy_to_modern[".xls"] == ".xlsx"
        assert converter._legacy_to_modern[".ppt"] == ".pptx"

    def test_legacy_extensions_set(self):
        """Legacy extensions are properly identified."""
        converter = OfficeConverter()
        assert ".doc" in converter._legacy_extensions
        assert ".xls" in converter._legacy_extensions
        assert ".ppt" in converter._legacy_extensions
        # Modern extensions are not legacy
        assert ".docx" not in converter._legacy_extensions

    def test_extract_from_legacy_removes_temp_file_and_directory(self, tmp_path, monkeypatch):
        """Legacy extraction cleans temporary LibreOffice artifacts after parsing."""
        converter = OfficeConverter()
        legacy_path = tmp_path / "legacy.doc"
        legacy_path.write_bytes(b"legacy-binary")

        temp_dir = tmp_path / "flavia_office_test"
        temp_dir.mkdir()
        converted_path = temp_dir / "legacy.docx"
        converted_path.write_bytes(b"converted")

        monkeypatch.setattr(
            converter,
            "_convert_with_libreoffice",
            lambda _path, _target_ext: converted_path,
        )
        monkeypatch.setattr(
            converter,
            "_extract_from_docx",
            lambda _path, is_odt=False: "converted text",
        )

        text = converter.extract_text(legacy_path)
        assert text == "converted text"
        assert not converted_path.exists()
        assert not temp_dir.exists()


class TestMarkdownFormatting:
    """Tests for markdown formatting utilities."""

    def test_strip_markdown(self):
        """_strip_markdown removes markdown formatting."""
        converter = OfficeConverter()

        # Test heading removal
        assert converter._strip_markdown("# Heading") == "Heading"
        assert converter._strip_markdown("## Heading") == "Heading"

        # Test bold removal
        assert converter._strip_markdown("**bold text**") == "bold text"

        # Test bullet removal
        assert converter._strip_markdown("- item") == "item"

        # Test blockquote removal
        assert converter._strip_markdown("> quoted") == "quoted"


class TestConversionPreservesStructure:
    """Tests for directory structure preservation."""

    @pytest.fixture
    def nested_docx(self, tmp_path):
        """Create a Word document in a nested directory."""
        docx = pytest.importorskip("docx")
        from docx import Document

        nested_dir = tmp_path / "docs" / "reports"
        nested_dir.mkdir(parents=True)

        doc = Document()
        doc.add_paragraph("Nested document content")
        doc_path = nested_dir / "report.docx"
        doc.save(doc_path)

        return doc_path

    def test_preserves_directory_structure(self, nested_docx, tmp_path):
        """Converted files preserve source directory structure."""
        converter = OfficeConverter()
        output_dir = tmp_path / ".converted"

        result = converter.convert(nested_docx, output_dir, output_format="md")

        assert result is not None
        assert result.exists()
        # The output should be in a matching nested structure
        assert "docs" in str(result) or "report.md" in str(result)
