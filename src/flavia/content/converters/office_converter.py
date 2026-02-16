"""Office document to text/markdown converter.

Supports Microsoft Office (.docx, .xlsx, .pptx) and legacy formats (.doc, .xls, .ppt)
as well as OpenDocument formats (.odt, .ods, .odp).
"""

import re
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Optional

from .base import BaseConverter


class OfficeConverter(BaseConverter):
    """Converts Office documents to text or markdown format."""

    supported_extensions = {
        # Modern Office
        ".docx",
        ".xlsx",
        ".pptx",
        # Legacy Office
        ".doc",
        ".xls",
        ".ppt",
        # OpenDocument
        ".odt",
        ".ods",
        ".odp",
    }

    requires_dependencies = ["python-docx", "openpyxl", "python-pptx"]
    dependency_import_map = {
        "python-docx": "docx",
        "openpyxl": "openpyxl",
        "python-pptx": "pptx",
    }

    # Extensions that require conversion to modern format first
    _legacy_extensions = {".doc", ".xls", ".ppt"}
    _legacy_to_modern = {
        ".doc": ".docx",
        ".xls": ".xlsx",
        ".ppt": ".pptx",
    }

    def convert(
        self,
        source_path: Path,
        output_dir: Path,
        output_format: str = "md",
    ) -> Optional[Path]:
        """Convert an Office document to text or markdown.

        Args:
            source_path: Source document path.
            output_dir: Output directory.
            output_format: "md" or "txt".

        Returns:
            Path to the converted file, or None on failure.
        """
        text = self.extract_text(source_path)
        if not text or not text.strip():
            return None

        if output_format == "md":
            content = text  # Already formatted as markdown
        else:
            # Strip markdown formatting for plain text
            content = self._strip_markdown(text)

        # Preserve directory structure when source lives under output_dir.parent.
        try:
            relative_source = source_path.resolve().relative_to(output_dir.resolve().parent)
            output_file = output_dir / relative_source.with_suffix(f".{output_format}")
        except ValueError:
            output_file = output_dir / (source_path.stem + f".{output_format}")

        output_file.parent.mkdir(parents=True, exist_ok=True)
        output_file.write_text(content, encoding="utf-8")
        return output_file

    def extract_text(self, source_path: Path) -> Optional[str]:
        """Extract text from an Office document.

        Args:
            source_path: Source document path.

        Returns:
            Extracted text as markdown, or None on failure.
        """
        ext = source_path.suffix.lower()

        # Handle legacy formats by converting to modern format first
        if ext in self._legacy_extensions:
            return self._extract_from_legacy(source_path)

        # Handle OpenDocument formats same as their modern Office equivalents
        if ext == ".odt":
            return self._extract_from_docx(source_path, is_odt=True)
        elif ext == ".ods":
            return self._extract_from_xlsx(source_path, is_ods=True)
        elif ext == ".odp":
            return self._extract_from_pptx(source_path, is_odp=True)

        # Modern Office formats
        if ext == ".docx":
            return self._extract_from_docx(source_path)
        elif ext == ".xlsx":
            return self._extract_from_xlsx(source_path)
        elif ext == ".pptx":
            return self._extract_from_pptx(source_path)

        return None

    def _extract_from_docx(self, path: Path, is_odt: bool = False) -> Optional[str]:
        """Extract text from Word document with heading and table detection.

        Args:
            path: Path to the document.
            is_odt: If True, convert from ODT format first.

        Returns:
            Markdown-formatted text, or None on failure.
        """
        try:
            from docx import Document
        except ImportError:
            return None

        try:
            if is_odt:
                # Convert ODT to DOCX first using LibreOffice
                converted = self._convert_with_libreoffice(path, ".docx")
                if not converted:
                    return None
                doc = Document(converted)
                # Clean up temp file after loading
                self._cleanup_temp_file(converted)
            else:
                doc = Document(path)
        except Exception:
            return None

        lines = []
        title = path.stem.replace("_", " ").replace("-", " ")
        lines.append(f"# {title}")
        lines.append("")

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # Detect heading styles
            style_name = para.style.name.lower() if para.style else ""
            if "heading 1" in style_name:
                lines.append(f"\n## {text}\n")
            elif "heading 2" in style_name:
                lines.append(f"\n### {text}\n")
            elif "heading 3" in style_name:
                lines.append(f"\n#### {text}\n")
            elif "title" in style_name:
                # Replace the auto-generated title with document title
                if lines and lines[0].startswith("# "):
                    lines[0] = f"# {text}"
                    lines.append("")
            else:
                lines.append(text)
                lines.append("")

        # Extract tables
        for table in doc.tables:
            table_md = self._table_to_markdown(table)
            if table_md:
                lines.append(table_md)
                lines.append("")

        return "\n".join(lines)

    def _extract_from_xlsx(self, path: Path, is_ods: bool = False) -> Optional[str]:
        """Extract text from Excel spreadsheet as markdown tables.

        Args:
            path: Path to the spreadsheet.
            is_ods: If True, convert from ODS format first.

        Returns:
            Markdown-formatted text, or None on failure.
        """
        try:
            from openpyxl import load_workbook
        except ImportError:
            return None

        try:
            if is_ods:
                # Convert ODS to XLSX first using LibreOffice
                converted = self._convert_with_libreoffice(path, ".xlsx")
                if not converted:
                    return None
                wb = load_workbook(converted, data_only=True)
                # Clean up temp file after loading
                self._cleanup_temp_file(converted)
            else:
                wb = load_workbook(path, data_only=True)
        except Exception:
            return None

        lines = []
        title = path.stem.replace("_", " ").replace("-", " ")
        lines.append(f"# {title}")
        lines.append("")

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            lines.append(f"## Sheet: {sheet_name}")
            lines.append("")

            # Get all rows with data
            rows = list(sheet.iter_rows(values_only=True))
            if not rows:
                lines.append("(empty sheet)")
                lines.append("")
                continue

            # Filter out completely empty rows
            rows = [row for row in rows if any(cell is not None for cell in row)]
            if not rows:
                lines.append("(empty sheet)")
                lines.append("")
                continue

            # Find the actual column width (ignore trailing None columns)
            max_cols = 0
            for row in rows:
                for i, cell in enumerate(reversed(row)):
                    if cell is not None:
                        max_cols = max(max_cols, len(row) - i)
                        break

            if max_cols == 0:
                lines.append("(empty sheet)")
                lines.append("")
                continue

            # Create markdown table
            # Header row
            header = rows[0][:max_cols]
            header_cells = [str(cell) if cell is not None else "" for cell in header]
            header_cells = [c.replace("|", "\\|") for c in header_cells]
            lines.append("| " + " | ".join(header_cells) + " |")
            lines.append("| " + " | ".join(["---"] * max_cols) + " |")

            # Data rows
            for row in rows[1:]:
                cells = row[:max_cols]
                cell_strs = [str(cell) if cell is not None else "" for cell in cells]
                # Escape pipe characters in cells
                cell_strs = [c.replace("|", "\\|") for c in cell_strs]
                lines.append("| " + " | ".join(cell_strs) + " |")

            lines.append("")

        return "\n".join(lines)

    def _extract_from_pptx(self, path: Path, is_odp: bool = False) -> Optional[str]:
        """Extract text from PowerPoint presentation.

        Args:
            path: Path to the presentation.
            is_odp: If True, convert from ODP format first.

        Returns:
            Markdown-formatted text, or None on failure.
        """
        try:
            from pptx import Presentation
        except ImportError:
            return None

        try:
            if is_odp:
                # Convert ODP to PPTX first using LibreOffice
                converted = self._convert_with_libreoffice(path, ".pptx")
                if not converted:
                    return None
                prs = Presentation(converted)
                # Clean up temp file after loading
                self._cleanup_temp_file(converted)
            else:
                prs = Presentation(path)
        except Exception:
            return None

        lines = []
        title = path.stem.replace("_", " ").replace("-", " ")
        lines.append(f"# {title}")
        lines.append("")

        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"## Slide {i}")
            lines.append("")

            # Extract title if present
            if slide.shapes.title and slide.shapes.title.text:
                lines.append(f"### {slide.shapes.title.text}")
                lines.append("")

            # Extract text from shapes
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                # Skip title shape as we already handled it
                if shape == slide.shapes.title:
                    continue

                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue

                    # Detect bullet points by indentation level
                    if para.level > 0:
                        indent = "  " * para.level
                        lines.append(f"{indent}- {text}")
                    elif any(
                        run.font.bold for run in para.runs if hasattr(run.font, "bold")
                    ):
                        # Bold text might be a sub-heading
                        lines.append(f"**{text}**")
                    else:
                        lines.append(f"- {text}")

            lines.append("")

            # Extract speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    lines.append("> **Notes:**")
                    for note_line in notes_text.split("\n"):
                        note_line = note_line.strip()
                        if note_line:
                            lines.append(f"> {note_line}")
                    lines.append("")

        return "\n".join(lines)

    def _extract_from_legacy(self, path: Path) -> Optional[str]:
        """Extract text from legacy Office formats by converting to modern format.

        Args:
            path: Path to the legacy document.

        Returns:
            Extracted text, or None on failure.
        """
        ext = path.suffix.lower()
        target_ext = self._legacy_to_modern.get(ext)
        if not target_ext:
            return None

        converted = self._convert_with_libreoffice(path, target_ext)
        if not converted:
            return None

        try:
            if target_ext == ".docx":
                return self._extract_from_docx(converted)
            elif target_ext == ".xlsx":
                return self._extract_from_xlsx(converted)
            elif target_ext == ".pptx":
                return self._extract_from_pptx(converted)
        finally:
            # Clean up temp file
            self._cleanup_temp_file(converted)

        return None

    def _convert_with_libreoffice(self, path: Path, target_ext: str) -> Optional[Path]:
        """Convert a document using LibreOffice CLI.

        Args:
            path: Path to the source document.
            target_ext: Target extension (e.g., ".docx").

        Returns:
            Path to converted file in temp directory, or None on failure.
        """
        # Check if LibreOffice is available
        libreoffice_cmd = self._find_libreoffice()
        if not libreoffice_cmd:
            return None

        # Determine output filter based on target format
        filter_map = {
            ".docx": "MS Word 2007 XML",
            ".xlsx": "Calc MS Excel 2007 XML",
            ".pptx": "Impress MS PowerPoint 2007 XML",
        }
        output_filter = filter_map.get(target_ext)
        if not output_filter:
            return None

        # Create temp directory for output
        temp_dir = Path(tempfile.mkdtemp(prefix="flavia_office_"))

        try:
            cmd = [
                libreoffice_cmd,
                "--headless",
                "--convert-to",
                target_ext.lstrip(".") + ":" + output_filter,
                "--outdir",
                str(temp_dir),
                str(path),
            ]

            result = subprocess.run(
                cmd,
                capture_output=True,
                timeout=60,
            )

            if result.returncode != 0:
                return None

            # Find the converted file
            converted_name = path.stem + target_ext
            converted_path = temp_dir / converted_name

            if not converted_path.exists():
                # Try to find any file with the target extension
                candidates = list(temp_dir.glob(f"*{target_ext}"))
                if candidates:
                    converted_path = candidates[0]
                else:
                    return None

            return converted_path

        except (subprocess.TimeoutExpired, subprocess.SubprocessError):
            return None

    def _find_libreoffice(self) -> Optional[str]:
        """Find LibreOffice executable.

        Returns:
            Path to LibreOffice command, or None if not found.
        """
        candidates = [
            "libreoffice",
            "soffice",
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            "/usr/bin/libreoffice",
            "/usr/local/bin/libreoffice",
        ]

        for cmd in candidates:
            if shutil.which(cmd):
                return cmd

        return None

    def _table_to_markdown(self, table) -> str:
        """Convert a Word table to markdown format.

        Args:
            table: python-docx Table object.

        Returns:
            Markdown table string.
        """
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
            rows.append(cells)

        if not rows:
            return ""

        # Calculate column widths
        num_cols = max(len(row) for row in rows) if rows else 0
        if num_cols == 0:
            return ""

        # Normalize rows to have same number of columns
        for row in rows:
            while len(row) < num_cols:
                row.append("")

        lines = []
        # Header row
        lines.append("| " + " | ".join(rows[0]) + " |")
        lines.append("| " + " | ".join(["---"] * num_cols) + " |")

        # Data rows
        for row in rows[1:]:
            lines.append("| " + " | ".join(row) + " |")

        return "\n".join(lines)

    @staticmethod
    def _strip_markdown(text: str) -> str:
        """Strip markdown formatting from text.

        Args:
            text: Markdown-formatted text.

        Returns:
            Plain text.
        """
        # Remove headings markers
        text = re.sub(r"^#+\s+", "", text, flags=re.MULTILINE)
        # Remove bold/italic
        text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
        text = re.sub(r"\*([^*]+)\*", r"\1", text)
        # Remove bullet points
        text = re.sub(r"^\s*[-*]\s+", "", text, flags=re.MULTILINE)
        # Remove blockquotes
        text = re.sub(r"^>\s*", "", text, flags=re.MULTILINE)
        # Remove table formatting
        text = re.sub(r"\|", " ", text)
        text = re.sub(r"-{3,}", "", text)
        return text

    @staticmethod
    def _cleanup_temp_file(path: Path) -> None:
        """Remove temporary converted file and best-effort cleanup of converter temp dir."""
        try:
            path.unlink()
        except Exception:
            pass

        try:
            parent = path.parent
            if parent.name.startswith("flavia_office_"):
                parent.rmdir()
        except Exception:
            pass
